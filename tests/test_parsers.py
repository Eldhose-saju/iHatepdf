"""One focused test per supported format parser (Task 3)."""
from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from conftest_helpers import ROOT as PROJECT_ROOT  # noqa: F401  (ensures src on sys.path)

from studymate.documents.parsers.docx import DOCXParser
from studymate.documents.parsers.image import ImageParser
from studymate.documents.parsers.pdf import PDFParser
from studymate.documents.parsers.pptx import PPTXParser


class TestDOCXParser(unittest.TestCase):
    def test_extracts_paragraph_text(self):
        import docx
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.docx"
            d = docx.Document()
            d.add_paragraph("Database systems use B-trees for indexing.")
            d.save(str(path))

            pages = DOCXParser().parse(path)
            self.assertEqual(len(pages), 1)
            self.assertIn("B-trees", pages[0].text)
            self.assertFalse(pages[0].is_empty)


class TestPPTXParser(unittest.TestCase):
    def test_extracts_one_page_per_slide(self):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.pptx"
            prs = Presentation()
            layout = prs.slide_layouts[1]
            for title in ("CPU Scheduling", "Memory Management"):
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = title
            prs.save(str(path))

            pages = PPTXParser().parse(path)
            self.assertEqual(len(pages), 2)
            self.assertIn("CPU Scheduling", pages[0].text)
            self.assertIn("Memory Management", pages[1].text)


class TestPDFParser(unittest.TestCase):
    def test_extracts_text_per_page(self):
        from reportlab.pdfgen import canvas
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.pdf"
            c = canvas.Canvas(str(path))
            c.drawString(72, 720, "Normalization removes redundant data.")
            c.showPage()
            c.drawString(72, 720, "A transaction must be atomic and durable.")
            c.showPage()
            c.save()

            pages = PDFParser().parse(path)
            self.assertEqual(len(pages), 2)
            self.assertIn("Normalization", pages[0].text)
            self.assertIn("atomic", pages[1].text)


class TestImageParser(unittest.TestCase):
    def test_returns_single_empty_page_for_ocr_service_to_fill(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "scan.png"
            from PIL import Image
            Image.new("RGB", (10, 10), color="white").save(path)

            pages = ImageParser().parse(path)
            self.assertEqual(len(pages), 1)
            self.assertTrue(pages[0].is_empty)


if __name__ == "__main__":
    unittest.main()
