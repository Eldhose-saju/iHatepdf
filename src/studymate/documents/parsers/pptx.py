"""PowerPoint extraction using python-pptx. One page per slide."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from studymate.documents.parsers.base import ParsedPage


class PPTXParser:
    def parse(self, filepath: Path) -> list[ParsedPage]:
        presentation = Presentation(str(filepath))
        pages: list[ParsedPage] = []
        for i, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                if getattr(shape, "has_notes_slide", False):
                    pass
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                parts.append(f"[notes] {slide.notes_slide.notes_text_frame.text}")
            pages.append(ParsedPage(page_number=i, text="\n".join(parts)))
        return pages
